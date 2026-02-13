from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x18, 0x18, 0x20)
PURPLE_PRIMARY = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
PURPLE_DARK = RGBColor(0x6D, 0x28, 0xD9)
TEXT_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_GRAY = RGBColor(0xA1, 0xA1, 0xAA)
TEXT_DIM = RGBColor(0x71, 0x71, 0x7A)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
PINK = RGBColor(0xEC, 0x48, 0x99)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)
YELLOW_ACCENT = RGBColor(0xFA, 0xCC, 0x15)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color, alpha=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        # python-pptx doesn't natively support alpha on shape fill easily,
        # so we just use a darker variant
        pass
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=TEXT_WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    """Add a text box with single-style text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(
    slide,
    left,
    top,
    width,
    height,
    items,
    font_size=16,
    color=TEXT_GRAY,
    bullet_color=PURPLE_PRIMARY,
    spacing=Pt(8),
):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25b8  "  # small triangle
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True
        run_bullet.font.name = "Calibri"
        # Text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Calibri"
    return txBox


def add_accent_bar(
    slide, left, top, width=Inches(0.06), height=Inches(0.8), color=PURPLE_PRIMARY
):
    """Add a small vertical accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_footer(slide, slide_num, total=4):
    """Add consistent footer with branding and slide number."""
    # Bottom bar
    bar = add_shape_fill(
        slide,
        Inches(0),
        SLIDE_HEIGHT - Inches(0.5),
        SLIDE_WIDTH,
        Inches(0.5),
        RGBColor(0x0F, 0x0F, 0x16),
    )

    # VH logo
    logo = add_rounded_rect(
        slide,
        Inches(0.5),
        SLIDE_HEIGHT - Inches(0.42),
        Inches(0.32),
        Inches(0.32),
        PURPLE_PRIMARY,
    )
    logo_tf = logo.text_frame
    logo_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = logo_tf.paragraphs[0].add_run()
    run.text = "VH"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = TEXT_WHITE
    run.font.name = "Calibri"

    add_text_box(
        slide,
        Inches(0.9),
        SLIDE_HEIGHT - Inches(0.44),
        Inches(2),
        Inches(0.35),
        "VibeHiring",
        font_size=11,
        color=TEXT_GRAY,
        bold=True,
    )

    add_text_box(
        slide,
        Inches(5.5),
        SLIDE_HEIGHT - Inches(0.44),
        Inches(3),
        Inches(0.35),
        "vibe-hiring-eight.vercel.app",
        font_size=10,
        color=TEXT_DIM,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        SLIDE_WIDTH - Inches(1.2),
        SLIDE_HEIGHT - Inches(0.44),
        Inches(1),
        Inches(0.35),
        f"{slide_num}/{total}",
        font_size=11,
        color=TEXT_DIM,
        alignment=PP_ALIGN.RIGHT,
    )


def create_slide_1_founder(prs):
    """Slide 1: Founder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)

    # Top-left badge
    badge = add_rounded_rect(
        slide,
        Inches(0.7),
        Inches(0.5),
        Inches(2.2),
        Inches(0.4),
        RGBColor(0x1A, 0x10, 0x2E),
    )
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "FOUNDER"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = PURPLE_LIGHT
    r.font.name = "Calibri"

    # Slide emoji + title
    add_text_box(
        slide,
        Inches(0.7),
        Inches(1.1),
        Inches(6),
        Inches(0.7),
        "\U0001f3a4  Meet the Founder",
        font_size=36,
        color=TEXT_WHITE,
        bold=True,
    )

    # ── Left column: Name + Vision ──
    # Name block
    add_accent_bar(slide, Inches(0.7), Inches(2.2), height=Inches(1.0))
    add_text_box(
        slide,
        Inches(1.0),
        Inches(2.2),
        Inches(5),
        Inches(0.5),
        "Karar",
        font_size=32,
        color=TEXT_WHITE,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(1.0),
        Inches(2.7),
        Inches(5),
        Inches(0.4),
        "Founder & Builder, VibeHiring",
        font_size=16,
        color=PURPLE_LIGHT,
    )

    # Vision box
    vision_box = add_rounded_rect(
        slide,
        Inches(0.7),
        Inches(3.5),
        Inches(5.5),
        Inches(1.6),
        RGBColor(0x14, 0x0E, 0x26),
    )
    vision_box.line.color.rgb = RGBColor(0x3F, 0x2A, 0x6E)
    vision_box.line.width = Pt(1)

    add_text_box(
        slide,
        Inches(1.0),
        Inches(3.65),
        Inches(5),
        Inches(0.35),
        "VISION",
        font_size=12,
        color=PURPLE_LIGHT,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(1.0),
        Inches(4.0),
        Inches(4.8),
        Inches(1.0),
        "Make hiring decisions based on evidence, not buzzwords.\n\nSaw firsthand how resumes fail to reflect actual engineering ability. Built VibeHiring to replicate senior-engineer judgment at scale.",
        font_size=14,
        color=TEXT_GRAY,
    )

    # ── Right column: Background + Why Me ──
    # Background section
    add_text_box(
        slide,
        Inches(7.0),
        Inches(2.0),
        Inches(5.5),
        Inches(0.4),
        "Background",
        font_size=20,
        color=TEXT_WHITE,
        bold=True,
    )

    add_bullet_list(
        slide,
        Inches(7.0),
        Inches(2.5),
        Inches(5.8),
        Inches(2.0),
        [
            "Builder focused on AI systems that interpret\nreal-world data, not keywords",
            "Experience designing agentic workflows\nand evaluation pipelines",
            "Obsessed with improving signal quality\nin technical hiring",
        ],
        font_size=14,
        spacing=Pt(12),
    )

    # Why Me section
    add_text_box(
        slide,
        Inches(7.0),
        Inches(4.5),
        Inches(5.5),
        Inches(0.4),
        "Why Me",
        font_size=20,
        color=TEXT_WHITE,
        bold=True,
    )

    add_bullet_list(
        slide,
        Inches(7.0),
        Inches(4.9),
        Inches(5.8),
        Inches(1.5),
        [
            "Saw firsthand how resumes fail to reflect\nactual engineering ability",
            "Built VibeHiring to replicate senior-engineer\njudgment at scale",
        ],
        font_size=14,
        spacing=Pt(12),
    )

    # Decorative gradient circle (top right)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10.5), Inches(-1), Inches(4), Inches(4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x14, 0x0E, 0x22)
    circle.line.fill.background()

    # Move circle to back
    sp = circle._element
    sp.getparent().insert(2, sp)

    add_footer(slide, 1)


def create_slide_2_problem(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Top badge
    badge = add_rounded_rect(
        slide,
        Inches(0.7),
        Inches(0.5),
        Inches(2.8),
        Inches(0.4),
        RGBColor(0x2A, 0x15, 0x05),
    )
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "\u26a0\ufe0f  THE PROBLEM"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = ORANGE
    r.font.name = "Calibri"

    # Title
    add_text_box(
        slide,
        Inches(0.7),
        Inches(1.2),
        Inches(10),
        Inches(0.8),
        "Tech hiring is noisy and inefficient",
        font_size=36,
        color=TEXT_WHITE,
        bold=True,
    )

    # ── Left column: Pain points ──
    pain_points = [
        ("Recruiters review hundreds of resumes manually", RED_ACCENT),
        ("ATS tools match keywords, not capability", RED_ACCENT),
        ("No reliable way to evaluate real coding ability", RED_ACCENT),
        ("GitHub portfolios are ignored or manually reviewed", RED_ACCENT),
        ("AI-generated code is rising \u2014 nobody measures authenticity", RED_ACCENT),
        ("Hiring decisions are opaque and hard to justify", RED_ACCENT),
    ]

    y_start = Inches(2.2)
    for i, (text, color) in enumerate(pain_points):
        y = y_start + Inches(i * 0.65)

        # X icon
        x_box = add_rounded_rect(
            slide,
            Inches(0.7),
            y,
            Inches(0.35),
            Inches(0.35),
            RGBColor(0x2A, 0x0A, 0x0A),
        )
        tf = x_box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = "\u2717"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = RED_ACCENT
        r.font.name = "Calibri"

        add_text_box(
            slide,
            Inches(1.2),
            y,
            Inches(5.5),
            Inches(0.4),
            text,
            font_size=15,
            color=TEXT_GRAY,
        )

    # ── Right column: Impact card ──
    impact_card = add_rounded_rect(
        slide,
        Inches(7.2),
        Inches(2.2),
        Inches(5.5),
        Inches(3.8),
        RGBColor(0x1A, 0x12, 0x08),
    )
    impact_card.line.color.rgb = RGBColor(0x4A, 0x2A, 0x0A)
    impact_card.line.width = Pt(1)

    add_text_box(
        slide,
        Inches(7.6),
        Inches(2.5),
        Inches(4.5),
        Inches(0.4),
        "THE COST",
        font_size=14,
        color=ORANGE,
        bold=True,
    )

    stats = [
        ("250+", "resumes per open role\non average"),
        ("23 hrs", "avg recruiter time per hire\non resume screening alone"),
        ("46%", "of new hires are considered\npoor fits within 18 months"),
    ]

    for i, (num, desc) in enumerate(stats):
        y = Inches(3.1) + Inches(i * 1.1)
        add_text_box(
            slide,
            Inches(7.6),
            y,
            Inches(1.5),
            Inches(0.5),
            num,
            font_size=28,
            color=ORANGE,
            bold=True,
        )
        add_text_box(
            slide,
            Inches(9.3),
            y,
            Inches(3.2),
            Inches(0.7),
            desc,
            font_size=13,
            color=TEXT_DIM,
        )

    # Bottom impact line
    impact_bar = add_rounded_rect(
        slide,
        Inches(0.7),
        Inches(6.2),
        Inches(12),
        Inches(0.6),
        RGBColor(0x14, 0x0A, 0x04),
    )
    impact_bar.line.color.rgb = RGBColor(0x3A, 0x20, 0x08)
    impact_bar.line.width = Pt(1)

    add_text_box(
        slide,
        Inches(1.0),
        Inches(6.25),
        Inches(11.5),
        Inches(0.45),
        "Result: Companies waste time, miss strong candidates, and hire with low confidence.",
        font_size=15,
        color=ORANGE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_footer(slide, 2)


def create_slide_3_solution(prs):
    """Slide 3: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Top badge
    badge = add_rounded_rect(
        slide,
        Inches(0.7),
        Inches(0.4),
        Inches(3.0),
        Inches(0.4),
        RGBColor(0x0A, 0x1A, 0x10),
    )
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "\U0001f680  THE SOLUTION"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = GREEN
    r.font.name = "Calibri"

    # Title
    add_text_box(
        slide,
        Inches(0.7),
        Inches(1.0),
        Inches(12),
        Inches(0.6),
        "VibeHiring: AI that evaluates like a senior engineer",
        font_size=32,
        color=TEXT_WHITE,
        bold=True,
    )

    # ── Feature cards in 2x3 grid ──
    features = [
        (
            "\U0001f4c4",
            "Multi-Format Parsing",
            "Reads PDF, Word, scanned images.\nOCR fallback for anything unreadable.",
            BLUE,
        ),
        (
            "\U0001f517",
            "Auto GitHub Discovery",
            "Extracts GitHub URLs from PDF\nhyperlinks automatically.",
            PINK,
        ),
        (
            "\U0001f50d",
            "Vibe Coding Detection",
            "Detects AI-generated vs human-written\ncode from GitHub repositories.",
            GREEN,
        ),
        (
            "\U0001f4ca",
            "6-Axis Radar Scoring",
            "Technical, Experience, Education,\nSoft Skills, Culture Fit, Growth.",
            PURPLE_PRIMARY,
        ),
        (
            "\U0001f4ac",
            "Explainable Rankings",
            '"Why This Rank?" narrative for\nevery candidate. No black boxes.',
            ORANGE,
        ),
        (
            "\U0001f4e5",
            "Bulk CSV Import",
            "Upload hundreds of candidates.\nFull AI analysis per row.",
            CYAN,
        ),
    ]

    for i, (icon, title, desc, color) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.7) + Inches(col * 4.2)
        y = Inches(1.85) + Inches(row * 1.65)

        card = add_rounded_rect(slide, x, y, Inches(3.9), Inches(1.45), BG_CARD)
        card.line.color.rgb = RGBColor(0x3F, 0x3F, 0x46)
        card.line.width = Pt(0.75)

        # Icon circle
        icon_bg = add_rounded_rect(
            slide,
            x + Inches(0.2),
            y + Inches(0.2),
            Inches(0.4),
            Inches(0.4),
            RGBColor(
                min(color[0] + 10, 255) if color[0] < 40 else color[0] // 5,
                min(color[1] + 10, 255) if color[1] < 40 else color[1] // 5,
                min(color[2] + 10, 255) if color[2] < 40 else color[2] // 5,
            ),
        )
        tf_i = icon_bg.text_frame
        tf_i.paragraphs[0].alignment = PP_ALIGN.CENTER
        r_i = tf_i.paragraphs[0].add_run()
        r_i.text = icon
        r_i.font.size = Pt(16)

        add_text_box(
            slide,
            x + Inches(0.75),
            y + Inches(0.18),
            Inches(3.0),
            Inches(0.35),
            title,
            font_size=15,
            color=TEXT_WHITE,
            bold=True,
        )

        add_text_box(
            slide,
            x + Inches(0.2),
            y + Inches(0.7),
            Inches(3.5),
            Inches(0.7),
            desc,
            font_size=12,
            color=TEXT_GRAY,
        )

    # ── Core Innovation box ──
    innov_box = add_rounded_rect(
        slide,
        Inches(0.7),
        Inches(5.3),
        Inches(6.0),
        Inches(1.2),
        RGBColor(0x14, 0x0E, 0x26),
    )
    innov_box.line.color.rgb = RGBColor(0x3F, 0x2A, 0x6E)
    innov_box.line.width = Pt(1)

    add_text_box(
        slide,
        Inches(1.0),
        Inches(5.4),
        Inches(5),
        Inches(0.3),
        "\u2699\ufe0f  CORE INNOVATION",
        font_size=12,
        color=PURPLE_LIGHT,
        bold=True,
    )

    add_text_box(
        slide,
        Inches(1.0),
        Inches(5.75),
        Inches(5.4),
        Inches(0.7),
        "Agentic AI that decides what evidence to gather\nbefore evaluating a candidate. Dual-model system:\nGemini 2.5 Flash (agent) + Gemini 2.5 Pro (evaluator).",
        font_size=13,
        color=TEXT_GRAY,
    )

    # ── Value box ──
    value_box = add_rounded_rect(
        slide,
        Inches(7.2),
        Inches(5.3),
        Inches(5.5),
        Inches(1.2),
        RGBColor(0x0A, 0x1A, 0x10),
    )
    value_box.line.color.rgb = RGBColor(0x15, 0x3A, 0x1E)
    value_box.line.width = Pt(1)

    add_text_box(
        slide,
        Inches(7.5),
        Inches(5.4),
        Inches(5),
        Inches(0.3),
        "\u2705  VALUE",
        font_size=12,
        color=GREEN,
        bold=True,
    )

    add_text_box(
        slide,
        Inches(7.5),
        Inches(5.75),
        Inches(5.0),
        Inches(0.7),
        "Better hiring decisions. Faster screening.\nDefensible rankings. Evidence over intuition.",
        font_size=14,
        color=TEXT_GRAY,
    )

    add_footer(slide, 3)


def create_slide_4_team(prs):
    """Slide 4: Team Needs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Top badge
    badge = add_rounded_rect(
        slide,
        Inches(0.7),
        Inches(0.5),
        Inches(2.8),
        Inches(0.4),
        RGBColor(0x0A, 0x12, 0x1E),
    )
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "\U0001f9e9  TEAM NEEDS"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = BLUE
    r.font.name = "Calibri"

    # Title
    add_text_box(
        slide,
        Inches(0.7),
        Inches(1.2),
        Inches(10),
        Inches(0.7),
        "To scale VibeHiring, we need",
        font_size=36,
        color=TEXT_WHITE,
        bold=True,
    )

    # ── Left: Technical roles ──
    tech_card = add_rounded_rect(
        slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(3.5), BG_CARD
    )
    tech_card.line.color.rgb = RGBColor(0x3F, 0x3F, 0x46)
    tech_card.line.width = Pt(1)

    add_accent_bar(
        slide,
        Inches(0.7),
        Inches(2.3),
        width=Inches(5.5),
        height=Inches(0.06),
        color=PURPLE_PRIMARY,
    )

    add_text_box(
        slide,
        Inches(1.1),
        Inches(2.55),
        Inches(4.5),
        Inches(0.4),
        "\u2699\ufe0f  Technical",
        font_size=22,
        color=TEXT_WHITE,
        bold=True,
    )

    # Role 1
    role1 = add_rounded_rect(
        slide,
        Inches(1.1),
        Inches(3.2),
        Inches(4.7),
        Inches(1.05),
        RGBColor(0x14, 0x0E, 0x26),
    )
    role1.line.color.rgb = RGBColor(0x3F, 0x2A, 0x6E)
    role1.line.width = Pt(0.75)

    add_text_box(
        slide,
        Inches(1.3),
        Inches(3.25),
        Inches(4.3),
        Inches(0.35),
        "ML Engineer",
        font_size=16,
        color=PURPLE_LIGHT,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(1.3),
        Inches(3.6),
        Inches(4.3),
        Inches(0.55),
        "Code analysis models, evaluation pipelines,\nvibe coding detection improvements",
        font_size=12,
        color=TEXT_GRAY,
    )

    # Role 2
    role2 = add_rounded_rect(
        slide,
        Inches(1.1),
        Inches(4.5),
        Inches(4.7),
        Inches(1.05),
        RGBColor(0x0A, 0x14, 0x1E),
    )
    role2.line.color.rgb = RGBColor(0x1A, 0x3A, 0x5A)
    role2.line.width = Pt(0.75)

    add_text_box(
        slide,
        Inches(1.3),
        Inches(4.55),
        Inches(4.3),
        Inches(0.35),
        "Full-Stack Engineer",
        font_size=16,
        color=BLUE,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(1.3),
        Inches(4.9),
        Inches(4.3),
        Inches(0.55),
        "Platform scaling, integrations, ATS connectors,\nreal-time collaboration features",
        font_size=12,
        color=TEXT_GRAY,
    )

    # ── Right: Business roles ──
    biz_card = add_rounded_rect(
        slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(3.5), BG_CARD
    )
    biz_card.line.color.rgb = RGBColor(0x3F, 0x3F, 0x46)
    biz_card.line.width = Pt(1)

    add_accent_bar(
        slide,
        Inches(7.0),
        Inches(2.3),
        width=Inches(5.5),
        height=Inches(0.06),
        color=GREEN,
    )

    add_text_box(
        slide,
        Inches(7.4),
        Inches(2.55),
        Inches(4.5),
        Inches(0.4),
        "\U0001f4bc  Business",
        font_size=22,
        color=TEXT_WHITE,
        bold=True,
    )

    # Role 3
    role3 = add_rounded_rect(
        slide,
        Inches(7.4),
        Inches(3.2),
        Inches(4.7),
        Inches(1.05),
        RGBColor(0x0A, 0x1A, 0x10),
    )
    role3.line.color.rgb = RGBColor(0x15, 0x3A, 0x1E)
    role3.line.width = Pt(0.75)

    add_text_box(
        slide,
        Inches(7.6),
        Inches(3.25),
        Inches(4.3),
        Inches(0.35),
        "B2B SaaS Growth / Sales Lead",
        font_size=16,
        color=GREEN,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(7.6),
        Inches(3.6),
        Inches(4.3),
        Inches(0.55),
        "Go-to-market strategy, enterprise sales,\nrecruiting agency partnerships",
        font_size=12,
        color=TEXT_GRAY,
    )

    # Role 4
    role4 = add_rounded_rect(
        slide,
        Inches(7.4),
        Inches(4.5),
        Inches(4.7),
        Inches(1.05),
        RGBColor(0x1A, 0x14, 0x08),
    )
    role4.line.color.rgb = RGBColor(0x4A, 0x3A, 0x1A)
    role4.line.width = Pt(0.75)

    add_text_box(
        slide,
        Inches(7.6),
        Inches(4.55),
        Inches(4.3),
        Inches(0.35),
        "Industry Advisor",
        font_size=16,
        color=YELLOW_ACCENT,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(7.6),
        Inches(4.9),
        Inches(4.3),
        Inches(0.55),
        "Recruiting or ATS industry experience,\nnetwork access, strategic guidance",
        font_size=12,
        color=TEXT_GRAY,
    )

    # ── Bottom: Goal banner ──
    goal_box = add_rounded_rect(
        slide,
        Inches(0.7),
        Inches(6.1),
        Inches(12),
        Inches(0.8),
        RGBColor(0x14, 0x0E, 0x26),
    )
    goal_box.line.color.rgb = RGBColor(0x3F, 0x2A, 0x6E)
    goal_box.line.width = Pt(1.5)

    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(6.2), Inches(11.5), Inches(0.55)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    r1 = p.add_run()
    r1.text = "GOAL:  "
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = PURPLE_LIGHT
    r1.font.name = "Calibri"

    r2 = p.add_run()
    r2.text = "Build the standard evaluation layer for technical hiring"
    r2.font.size = Pt(16)
    r2.font.bold = True
    r2.font.color.rgb = TEXT_WHITE
    r2.font.name = "Calibri"

    add_footer(slide, 4)


def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = Emu(12192000)  # 13.333 inches
    prs.slide_height = Emu(6858000)  # 7.5 inches

    create_slide_1_founder(prs)
    create_slide_2_problem(prs)
    create_slide_3_solution(prs)
    create_slide_4_team(prs)

    output = "VibeHiring_Pitch.pptx"
    prs.save(output)
    print(f"Created: {output}")


if __name__ == "__main__":
    main()
